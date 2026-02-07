#!/usr/bin/env python3
"""
Generate PPT/PPTX fixture pairs + screenshots using:
- python-pptx for authoring
- LibreOffice (soffice) for PPT/PPTX/PDF conversion
- pdftoppm for PDF -> PNG

This script is meant for maintainers of the fixture corpus.
"""

from __future__ import annotations

import os, json, glob, shutil, subprocess, datetime
from dataclasses import dataclass
from typing import List, Callable, Dict, Any, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

from PIL import Image, ImageDraw


ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
CASES_DIR = os.path.join(ROOT, "cases")
ASSETS_DIR = os.path.join(ROOT, "assets")
LO_PROFILE = os.path.join("/tmp", "lo-profile-webpptx-fixtures")


def run(cmd: List[str]) -> Tuple[int, str, str]:
    res = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    return res.returncode, res.stdout, res.stderr


def soffice_convert(input_path: str, outdir: str, out_ext: str) -> str:
    os.makedirs(outdir, exist_ok=True)
    cmd = [
        "soffice",
        f"-env:UserInstallation=file://{LO_PROFILE}",
        "--headless",
        "--nologo",
        "--convert-to",
        out_ext,
        "--outdir",
        outdir,
        input_path,
    ]
    rc, out, err = run(cmd)
    if rc != 0:
        raise RuntimeError(f"soffice failed ({rc})\nstdout:\n{out}\nstderr:\n{err}\ncmd:{cmd}")
    base = os.path.splitext(os.path.basename(input_path))[0]
    candidates = glob.glob(os.path.join(outdir, base + "." + out_ext))
    if candidates:
        return candidates[0]
    cand2 = sorted(glob.glob(os.path.join(outdir, f"*.{out_ext}")), key=os.path.getmtime, reverse=True)
    if cand2:
        return cand2[0]
    raise FileNotFoundError(f"Converted file not found: {input_path} -> {out_ext} in {outdir}")


def pdftoppm_png(pdf_path: str, outprefix: str, dpi: int = 72) -> List[str]:
    os.makedirs(os.path.dirname(outprefix), exist_ok=True)
    cmd = ["pdftoppm", "-png", "-r", str(dpi), pdf_path, outprefix]
    rc, out, err = run(cmd)
    if rc != 0:
        raise RuntimeError(f"pdftoppm failed ({rc})\nstdout:{out}\nstderr:{err}")
    return sorted(glob.glob(outprefix + "-*.png"))


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    return prs


def set_run_style(run, **kw):
    f = run.font
    if "size" in kw: f.size = Pt(kw["size"])
    if "bold" in kw: f.bold = kw["bold"]
    if "italic" in kw: f.italic = kw["italic"]
    if "underline" in kw: f.underline = kw["underline"]
    if "color" in kw: f.color.rgb = RGBColor(*kw["color"])


def add_title(slide, text: str):
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(0.3), Inches(12), Inches(0.8))
    tf = tx.text_frame
    tf.text = text
    run = tf.paragraphs[0].runs[0]
    set_run_style(run, size=40, bold=True, color=(0, 0, 0))
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    return tx


def build_assets() -> Tuple[str, str]:
    os.makedirs(ASSETS_DIR, exist_ok=True)
    png_path = os.path.join(ASSETS_DIR, "sample.png")
    jpg_path = os.path.join(ASSETS_DIR, "sample.jpg")

    img = Image.new("RGBA", (800, 450), (255, 255, 255, 255))
    d = ImageDraw.Draw(img)
    d.rectangle([50, 50, 350, 200], fill=(255, 99, 71, 255), outline=(0, 0, 0, 255), width=4)
    d.ellipse([420, 60, 740, 260], fill=(65, 105, 225, 255), outline=(0, 0, 0, 255), width=4)
    d.line([80, 350, 720, 350], fill=(34, 139, 34, 255), width=12)
    for y in range(280, 440, 20):
        for x in range(50, 250, 20):
            if (x // 20 + y // 20) % 2 == 0:
                d.rectangle([x, y, x + 20, y + 20], fill=(220, 220, 220, 255))
    img.save(png_path, "PNG")
    img.convert("RGB").save(jpg_path, "JPEG", quality=90)
    return png_path, jpg_path


@dataclass
class CaseSpec:
    id: str
    title: str
    description: str
    tags: List[str]
    slides: int
    level: str
    build: Callable[[Presentation, str, str], None]


def cases(png_path: str, jpg_path: str) -> List[CaseSpec]:
    def blank(prs: Presentation, *_):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 245, 245)
        tx = s.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(12), Inches(0.8))
        tf = tx.text_frame
        tf.text = "blank baseline"
        set_run_style(tf.paragraphs[0].runs[0], size=32, bold=True, color=(120, 120, 120))
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def text_basic(prs: Presentation, *_):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(s, "Text: runs & styles")
        box = s.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(12.0), Inches(4.8))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.text = ""
        r1 = p.add_run(); r1.text = "Bold "; set_run_style(r1, size=28, bold=True, color=(0, 0, 0))
        r2 = p.add_run(); r2.text = "Italic "; set_run_style(r2, size=28, italic=True, color=(0, 0, 0))
        r3 = p.add_run(); r3.text = "Underline "; set_run_style(r3, size=28, underline=True, color=(0, 0, 0))
        r4 = p.add_run(); r4.text = "Color"; set_run_style(r4, size=28, color=(220, 20, 60))
        p2 = tf.add_paragraph()
        p2.text = ""
        for t, size in [("12 ", 12), ("18 ", 18), ("24 ", 24), ("36", 36)]:
            rr = p2.add_run(); rr.text = t; set_run_style(rr, size=size, color=(0, 100, 0))
        p3 = tf.add_paragraph()
        p3.text = "日本語テキスト / Emoji 😀"
        if p3.runs:
            set_run_style(p3.runs[0], size=24, color=(0, 0, 128))

    def bullets(prs: Presentation, *_):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(s, "Text: bullets & indent")
        tb = s.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(12), Inches(5.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.clear()
        items = [
            ("Level 0 - item A", 0),
            ("Level 0 - item B", 0),
            ("Level 1 - sub 1", 1),
            ("Level 1 - sub 2", 1),
            ("Level 2 - subsub", 2),
            ("Back to level 0", 0),
        ]
        for i, (t, lev) in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = t
            p.level = lev
            p.font.size = Pt(24 if lev == 0 else 20)

    def shapes(prs: Presentation, *_):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(s, "Shapes: basic")
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.7), Inches(4.0), Inches(2.0))
        rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(255, 215, 0)
        rect.line.color.rgb = RGBColor(0, 0, 0); rect.line.width = Pt(2)
        rect.text = "Rectangle"; rect.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        ell = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.0), Inches(1.7), Inches(4.0), Inches(2.0))
        ell.fill.solid(); ell.fill.fore_color.rgb = RGBColor(135, 206, 250)
        ell.line.color.rgb = RGBColor(25, 25, 112); ell.line.width = Pt(3)
        ell.text = "Oval"; ell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def line_styles(prs: Presentation, *_):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(s, "Lines: dash & round-dot")
        l1 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(2), Inches(12), Inches(2))
        l1.line.color.rgb = RGBColor(0, 0, 0); l1.line.width = Pt(3)
        l2 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(3.2), Inches(12), Inches(3.2))
        l2.line.color.rgb = RGBColor(220, 20, 60); l2.line.width = Pt(4)
        l2.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        l3 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(4.4), Inches(12), Inches(4.4))
        l3.line.color.rgb = RGBColor(30, 144, 255); l3.line.width = Pt(4)
        l3.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT

    def images(prs: Presentation, *_):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(s, "Images: PNG & JPG")
        s.shapes.add_picture(png_path, Inches(0.8), Inches(1.7), width=Inches(5.8))
        s.shapes.add_picture(jpg_path, Inches(7.0), Inches(1.7), width=Inches(5.8))

    def crop_rotate(prs: Presentation, *_):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(s, "Images: crop & rotate")
        pic = s.shapes.add_picture(png_path, Inches(1.2), Inches(1.8), width=Inches(5.5))
        pic.crop_left, pic.crop_right, pic.crop_top, pic.crop_bottom = 0.15, 0.15, 0.10, 0.10
        pic2 = s.shapes.add_picture(png_path, Inches(7.0), Inches(2.0), width=Inches(5.0))
        pic2.rotation = -12

    def table_basic(prs: Presentation, *_):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(s, "Table: basic")
        rows, cols = 4, 4
        ts = s.shapes.add_table(rows, cols, Inches(1.0), Inches(1.8), Inches(11.5), Inches(4.8))
        t = ts.table
        for j, h in enumerate(["A", "B", "C", "D"]):
            cell = t.cell(0, j)
            cell.text = h
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(47, 79, 79)

    def chart_bar(prs: Presentation, *_):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(s, "Chart: clustered column")
        cd = ChartData()
        cd.categories = ["A", "B", "C", "D"]
        cd.add_series("Series 1", (19.2, 21.4, 16.7, 23.2))
        cd.add_series("Series 2", (12.1, 14.2, 15.6, 10.9))
        s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.0), Inches(1.8), Inches(11.5), Inches(5.2), cd)

    def hyperlink(prs: Presentation, *_):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(s, "Hyperlink")
        tb = s.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.5), Inches(2.0))
        tf = tb.text_frame
        tf.text = "Open LibreOffice website"
        run = tf.paragraphs[0].runs[0]
        set_run_style(run, size=32, underline=True, color=(0, 0, 255))
        run.hyperlink.address = "https://www.libreoffice.org/"

    def multislide(prs: Presentation, *_):
        s1 = prs.slides.add_slide(prs.slide_layouts[6]); add_title(s1, "Multi-slide: 1/2")
        s2 = prs.slides.add_slide(prs.slide_layouts[6]); add_title(s2, "Multi-slide: 2/2")
        s2.shapes.add_picture(jpg_path, Inches(2.0), Inches(2.0), width=Inches(9.3))

    return [
        CaseSpec("000_blank", "Blank baseline", "Single blank slide with background fill and centered text.", ["slide", "background", "text"], 1, "core", blank),
        CaseSpec("010_text_basic", "Text basic", "Rich text runs: bold/italic/underline, colors, mixed font sizes, Japanese + emoji.", ["text", "runs", "unicode"], 1, "core", text_basic),
        CaseSpec("011_text_bullets", "Text bullets", "Bulleted list with multiple nesting levels.", ["text", "bullets"], 1, "core", bullets),
        CaseSpec("020_shapes_basic", "Shapes basic", "Rectangle/oval with fills and lines.", ["shapes", "fill", "line"], 1, "core", shapes),
        CaseSpec("021_line_styles", "Line styles", "Straight connectors with solid/dash/round-dot styles.", ["lines", "dash"], 1, "extended", line_styles),
        CaseSpec("030_images_basic", "Images basic", "PNG and JPG pictures placed/scaled.", ["images"], 1, "core", images),
        CaseSpec("031_images_crop_rotate", "Images crop & rotate", "Picture crop and rotation.", ["images", "crop", "rotation"], 1, "extended", crop_rotate),
        CaseSpec("040_table_basic", "Table basic", "Simple table.", ["table"], 1, "extended", table_basic),
        CaseSpec("050_chart_bar", "Chart bar", "Clustered column chart.", ["chart"], 1, "hard", chart_bar),
        CaseSpec("060_hyperlink", "Hyperlink", "Text run hyperlink.", ["hyperlink"], 1, "extended", hyperlink),
        CaseSpec("080_multislide", "Multi-slide", "Two slides.", ["slides"], 2, "core", multislide),
    ]


def render_to_pngs(in_path: str, fmt_dir: str, tmp: str, dpi: int = 72):
    os.makedirs(fmt_dir, exist_ok=True)
    for f in glob.glob(os.path.join(fmt_dir, "slide-*.png")):
        os.remove(f)
    pdf_out = soffice_convert(in_path, tmp, "pdf")
    pdf_target = os.path.join(fmt_dir, "slides.pdf")
    if os.path.exists(pdf_target):
        os.remove(pdf_target)
    shutil.move(pdf_out, pdf_target)

    raw_prefix = os.path.join(fmt_dir, "__raw_slide")
    for f in glob.glob(raw_prefix + "-*.png"):
        os.remove(f)
    pngs = pdftoppm_png(pdf_target, raw_prefix, dpi=dpi)
    for idx, fp in enumerate(pngs, start=1):
        os.replace(fp, os.path.join(fmt_dir, f"slide-{idx:03d}.png"))
    for f in glob.glob(raw_prefix + "-*.png"):
        os.remove(f)


def main():
    os.makedirs(CASES_DIR, exist_ok=True)
    os.makedirs(LO_PROFILE, exist_ok=True)

    png_path, jpg_path = build_assets()
    specs = cases(png_path, jpg_path)

    # build author decks
    for c in specs:
        cdir = os.path.join(CASES_DIR, c.id)
        os.makedirs(cdir, exist_ok=True)
        prs = new_prs()
        c.build(prs, png_path, jpg_path)
        author = os.path.join(cdir, "author.pptx")
        prs.save(author)

        meta = {
            "id": c.id,
            "title": c.title,
            "description": c.description,
            "tags": c.tags,
            "slides": c.slides,
            "level": c.level,
            "generated_at": datetime.datetime.utcnow().isoformat() + "Z",
        }
        with open(os.path.join(cdir, "meta.json"), "w", encoding="utf-8") as f:
            json.dump(meta, f, ensure_ascii=False, indent=2)

    # convert + render
    for c in specs:
        cdir = os.path.join(CASES_DIR, c.id)
        tmp = os.path.join(cdir, "_lo_out")
        if os.path.exists(tmp):
            shutil.rmtree(tmp)
        os.makedirs(tmp, exist_ok=True)

        author = os.path.join(cdir, "author.pptx")
        ref_ppt = os.path.join(cdir, "ref.ppt")
        ref_pptx = os.path.join(cdir, "ref.pptx")

        # author.pptx -> ref.ppt
        out_ppt = soffice_convert(author, tmp, "ppt")
        if os.path.exists(ref_ppt):
            os.remove(ref_ppt)
        shutil.move(out_ppt, ref_ppt)

        # ref.ppt -> ref.pptx
        out_pptx = soffice_convert(ref_ppt, tmp, "pptx")
        if os.path.exists(ref_pptx):
            os.remove(ref_pptx)
        shutil.move(out_pptx, ref_pptx)

        # render both
        render_dir = os.path.join(cdir, "render")
        render_to_pngs(ref_ppt, os.path.join(render_dir, "ppt"), tmp)
        render_to_pngs(ref_pptx, os.path.join(render_dir, "pptx"), tmp)

        shutil.rmtree(tmp, ignore_errors=True)

    # write manifest
    manifest = {
        "generated_at": datetime.datetime.utcnow().isoformat() + "Z",
        "slide_size": {"width_in": 13.333, "height_in": 7.5, "pixel_at_72dpi": [960, 540]},
        "render_pipeline": ["libreoffice:convert-to pdf", "pdftoppm -png -r 72"],
        "cases": [
            {
                "id": c.id,
                "title": c.title,
                "description": c.description,
                "tags": c.tags,
                "slides": c.slides,
                "level": c.level,
            }
            for c in specs
        ],
    }
    with open(os.path.join(ROOT, "manifest.json"), "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)

    print("done")


if __name__ == "__main__":
    main()
